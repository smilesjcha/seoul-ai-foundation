"""Slide layout helpers — see docs/09-ppt-design-system.md."""

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import theme as T
import shapes as S


def _add_blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = T.WHITE
    return slide


def _add_title(slide, text, *, top=None, left=None, width=None, color=None,
               size=None):
    return S.add_textbox(
        slide,
        left or T.CONTENT_LEFT,
        top or T.CONTENT_TOP,
        width or T.CONTENT_WIDTH,
        Inches(0.7),
        text,
        size=size or T.PT_TITLE,
        bold=True,
        color=color or T.BLACK,
    )


def _add_footer(slide, slide_no, total=77):
    S.add_textbox(slide, Inches(0.6), T.SLIDE_H - Inches(0.4), Inches(6),
                  Inches(0.3), "개발자에서 AI Native로 · 서울AI재단 5시간 워크숍",
                  size=T.PT_CAPTION, color=T.GRAY_80)
    S.add_textbox(slide, T.SLIDE_W - Inches(1.2), T.SLIDE_H - Inches(0.4),
                  Inches(0.6), Inches(0.3),
                  f"{slide_no:02d} / {total:02d}",
                  size=T.PT_CAPTION, color=T.GRAY_80, align=PP_ALIGN.RIGHT)


def _attach_note(slide, data):
    parts = []
    if data.get("note"):
        parts.append(data["note"].strip())
    if data.get("speak_seconds"):
        parts.append(f"[권장 발화 시간: {data['speak_seconds']}]")
    S.set_slide_notes(slide, "\n\n".join(parts))


# ---------------------------------------------------------------- Cover

def make_cover(prs, data):
    slide = _add_blank_slide(prs)
    S.add_textbox(slide, Inches(0.7), Inches(0.55), Inches(8), Inches(0.4),
                  "서울AI재단 5시간 워크숍",
                  size=T.PT_CAPTION, color=T.DEEP_BLUE, bold=True)
    S.add_thin_line(slide, Inches(0.7), Inches(1.0), Inches(1.6),
                    color=T.DEEP_BLUE, weight=Pt(2))
    title_box = S.add_textbox(slide, Inches(0.7), Inches(2.2), Inches(12),
                              Inches(1.3), " ", size=T.PT_COVER_TITLE,
                              bold=True, color=T.BLACK)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = ""
    r1 = p.add_run(); r1.text = "개발자에서 "
    S.set_run_font(r1, size=T.PT_COVER_TITLE, bold=True, color=T.BLACK)
    r2 = p.add_run(); r2.text = "AI Native"
    S.set_run_font(r2, size=T.PT_COVER_TITLE, bold=True, color=T.DEEP_BLUE)
    r3 = p.add_run(); r3.text = "로"
    S.set_run_font(r3, size=T.PT_COVER_TITLE, bold=True, color=T.BLACK)
    S.add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12),
                  Inches(0.6),
                  "Claude Desktop으로 완성하는 업무 생산성 혁신",
                  size=T.PT_COVER_SUB, color=T.GRAY_80)
    body = data.get("body", [])
    if body:
        S.add_multi_line_textbox(
            slide, Inches(0.7), Inches(4.4), Inches(12), Inches(1.6),
            body, size=T.PT_BODY, color=T.BLACK, line_spacing=1.4,
            bullet=True,
        )
    S.add_thin_line(slide, Inches(0.7), Inches(6.4), Inches(12.0),
                    color=T.GRAY_20, weight=Pt(0.5))
    S.add_textbox(slide, Inches(0.7), Inches(6.55), Inches(12),
                  Inches(0.4),
                  "강사 · 무신사 Agentic AI PM · 서울시립대 / 아주대 AI 부문 겸임교수",
                  size=T.PT_BODY_SM, color=T.GRAY_80)
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Section

def make_section_divider(prs, data):
    slide = _add_blank_slide(prs)
    part_no = data.get("part_no", "")
    S.add_textbox(slide, Inches(0.7), Inches(1.4), Inches(5),
                  Inches(0.5),
                  f"PART {part_no}",
                  size=T.PT_CAPTION, color=T.DEEP_BLUE, bold=True)
    S.add_textbox(slide, Inches(0.55), Inches(1.8), Inches(5.5),
                  Inches(3.5),
                  part_no, size=T.PT_SECTION_NUM, bold=True,
                  color=T.DEEP_BLUE)
    S.add_thin_line(slide, Inches(6.3), Inches(2.4), Inches(0.0),
                    color=T.DEEP_BLUE)
    rect = S.add_rect(slide, Inches(6.3), Inches(2.4), Emu(45720),
                      Inches(2.5), fill=T.DEEP_BLUE)
    title = data.get("section_title") or data.get("title", "")
    S.add_textbox(slide, Inches(6.6), Inches(2.4), Inches(6.2),
                  Inches(0.8), title,
                  size=T.PT_SECTION_TITLE, bold=True, color=T.BLACK,
                  line_spacing=1.15)
    sub = data.get("section_sub", "")
    if sub:
        S.add_textbox(slide, Inches(6.6), Inches(3.6), Inches(6.2),
                      Inches(0.5), sub,
                      size=T.PT_LEAD, color=T.GRAY_80)
    outcomes = data.get("section_outcomes", [])
    if outcomes:
        S.add_textbox(slide, Inches(6.6), Inches(4.3), Inches(6.2),
                      Inches(0.4), "이 파트의 산출물",
                      size=T.PT_CAPTION, color=T.DEEP_BLUE, bold=True)
        S.add_multi_line_textbox(slide, Inches(6.6), Inches(4.7),
                                 Inches(6.2), Inches(2.0),
                                 outcomes, size=T.PT_BODY, color=T.BLACK,
                                 line_spacing=1.3, bullet=True)
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Concept

def make_concept(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    body = data.get("body", [])
    lead = data.get("lead")
    body_top = Inches(1.9)
    if lead:
        S.add_textbox(slide, T.CONTENT_LEFT, Inches(1.7),
                      T.CONTENT_WIDTH, Inches(1.0),
                      lead, size=T.PT_LEAD, color=T.DEEP_BLUE,
                      bold=True, line_spacing=1.3)
        body_top = Inches(2.9)
    if body:
        S.add_multi_line_textbox(
            slide, T.CONTENT_LEFT, body_top,
            T.CONTENT_WIDTH, Inches(4.5),
            body, size=T.PT_BODY, color=T.BLACK,
            line_spacing=1.55, bullet=True,
        )
    visual_hint = data.get("visual_hint")
    if visual_hint:
        S.add_textbox(slide, T.CONTENT_LEFT, Inches(6.5),
                      T.CONTENT_WIDTH, Inches(0.4),
                      f"도식 · {visual_hint}",
                      size=T.PT_CAPTION, color=T.GRAY_80)
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Process

def make_process(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    items = data.get("flow_items") or []
    flow_top = Inches(2.4)
    flow_height = Inches(1.4)
    if items:
        S.add_horizontal_flow(slide, T.CONTENT_LEFT, flow_top,
                              T.CONTENT_WIDTH, flow_height, items,
                              emphasize_index=data.get("emphasize_index"))
    body = data.get("body", [])
    if body:
        S.add_multi_line_textbox(
            slide, T.CONTENT_LEFT, Inches(4.4),
            T.CONTENT_WIDTH, Inches(2.6),
            body, size=T.PT_BODY, color=T.BLACK,
            line_spacing=1.5, bullet=True,
        )
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Two-column

def make_two_column(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    col_w = (T.CONTENT_WIDTH - Inches(0.4)) / 2
    body = data.get("body", [])
    S.add_multi_line_textbox(
        slide, T.CONTENT_LEFT, Inches(1.9),
        col_w, Inches(5.0),
        body, size=T.PT_BODY, color=T.BLACK,
        line_spacing=1.55, bullet=True,
    )
    right_left = T.CONTENT_LEFT + col_w + Inches(0.4)
    right_kind = data.get("right_kind", "placeholder")
    if right_kind == "code":
        S.add_code_block(slide, right_left, Inches(1.9), col_w, Inches(4.8),
                         data.get("right_lines", []),
                         highlight_indices=data.get("right_highlight"))
    elif right_kind == "doc_box":
        rect_left = right_left
        rect_w = col_w
        S.add_filled_text_rect(
            slide, rect_left, Inches(1.9), rect_w, Inches(0.55),
            data.get("right_box_title", "문서 예시"),
            fill=T.DEEP_BLUE, line=None, text_color=T.WHITE,
            size=T.PT_BODY_SM, bold=True, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.MIDDLE)
        S.add_rect(slide, rect_left, Inches(2.45), rect_w,
                   Inches(4.3), fill=T.BLUE_10, line=T.DEEP_BLUE,
                   line_width=Pt(0.75))
        pad = Inches(0.2)
        S.add_multi_line_textbox(
            slide, rect_left + pad, Inches(2.45) + pad,
            rect_w - 2 * pad, Inches(4.0),
            data.get("right_box_lines", []), size=T.PT_BODY_SM,
            color=T.BLACK, line_spacing=1.4, bullet=True)
    elif right_kind == "cards":
        cards = data.get("right_cards", [])
        n = len(cards)
        if n:
            card_h = (Inches(5.0) - Inches(0.2) * (n - 1)) / n
            for i, c in enumerate(cards):
                top = Inches(1.9) + (card_h + Inches(0.2)) * i
                S.add_card(slide, right_left, top, col_w, card_h,
                           c["title"], c.get("body", []),
                           accent=c.get("accent", False))
    elif right_kind == "checklist":
        S.add_check_list(slide, right_left, Inches(1.9), col_w, Inches(5.0),
                         data.get("right_items", []), size=T.PT_BODY)
    else:
        S.add_placeholder_box(slide, right_left, Inches(1.9), col_w,
                              Inches(4.8),
                              data.get("right_caption",
                                       "도식 placeholder"))
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Document

def make_document(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    body = data.get("body", [])
    badge_text = data.get("badge")
    body_top = Inches(1.9)
    if badge_text:
        S.add_badge(slide, T.CONTENT_LEFT, Inches(1.85), Inches(1.6),
                    Inches(0.42), badge_text)
        body_top = Inches(2.5)
    S.add_multi_line_textbox(
        slide, T.CONTENT_LEFT, body_top,
        Inches(6.0), Inches(5.0),
        body, size=T.PT_BODY, color=T.BLACK,
        line_spacing=1.55, bullet=True,
    )
    box_left = T.CONTENT_LEFT + Inches(6.7)
    box_w = T.CONTENT_WIDTH - Inches(6.7)
    box_top = Inches(1.9)
    box_h = Inches(5.0)
    box_title = data.get("box_title", "문서 예시")
    S.add_filled_text_rect(
        slide, box_left, box_top, box_w, Inches(0.55), box_title,
        fill=T.DEEP_BLUE, line=None, text_color=T.WHITE,
        size=T.PT_BODY_SM, bold=True, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.MIDDLE)
    S.add_rect(slide, box_left, box_top + Inches(0.55), box_w,
               box_h - Inches(0.55), fill=T.BLUE_10, line=T.DEEP_BLUE,
               line_width=Pt(0.75))
    pad = Inches(0.25)
    box_lines = data.get("box_lines", [])
    S.add_multi_line_textbox(
        slide, box_left + pad, box_top + Inches(0.55) + pad,
        box_w - 2 * pad, box_h - Inches(0.55) - 2 * pad,
        box_lines, size=T.PT_BODY_SM, color=T.BLACK,
        line_spacing=1.4, bullet=False,
    )
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Code Demo

def make_code_demo(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    body = data.get("body", [])
    S.add_multi_line_textbox(
        slide, T.CONTENT_LEFT, Inches(1.9),
        Inches(5.4), Inches(5.0),
        body, size=T.PT_BODY, color=T.BLACK,
        line_spacing=1.55, bullet=True,
    )
    code_left = T.CONTENT_LEFT + Inches(6.0)
    code_w = T.CONTENT_WIDTH - Inches(6.0)
    S.add_code_block(slide, code_left, Inches(1.9), code_w, Inches(5.0),
                     data.get("code_lines", []),
                     highlight_indices=data.get("code_highlight"))
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Review

def make_review(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    lead = data.get("lead")
    items_top = Inches(1.9)
    if lead:
        S.add_callout_box(slide, T.CONTENT_LEFT, Inches(1.7),
                          T.CONTENT_WIDTH, Inches(0.8), lead,
                          size=T.PT_LEAD, bold=True)
        items_top = Inches(2.8)
    items = data.get("body", [])
    S.add_check_list(slide, T.CONTENT_LEFT + Inches(0.2), items_top,
                     T.CONTENT_WIDTH - Inches(0.4), Inches(4.0),
                     items, size=T.PT_BODY, line_spacing=1.7)
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Wrap-up

def make_wrap_up(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    body = data.get("body", [])
    S.add_multi_line_textbox(
        slide, T.CONTENT_LEFT, Inches(2.0),
        T.CONTENT_WIDTH, Inches(4.5),
        body, size=Pt(20), color=T.BLACK, bold=True,
        line_spacing=1.7, bullet=True,
    )
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Quote / One-liner

def make_one_liner(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    msg = data.get("body", [""])[0] if data.get("body") else ""
    bar = S.add_rect(slide, T.CONTENT_LEFT, Inches(2.6), Emu(45720),
                     Inches(2.5), fill=T.DEEP_BLUE)
    S.add_textbox(slide, T.CONTENT_LEFT + Inches(0.3), Inches(2.6),
                  T.CONTENT_WIDTH - Inches(0.3), Inches(2.5), msg,
                  size=Pt(28), color=T.BLACK, bold=True,
                  anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Custom: schedule table

def make_schedule(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    rows = data.get("schedule_rows", [])
    table = S.add_table(slide, T.CONTENT_LEFT, Inches(1.85),
                        T.CONTENT_WIDTH, Inches(4.8),
                        ["시간", "세션", "핵심 내용", "산출물"], rows,
                        col_widths=[Inches(1.7), Inches(2.6), Inches(5.7),
                                    Inches(2.1)])
    note = data.get("body", [])
    if note:
        S.add_multi_line_textbox(slide, T.CONTENT_LEFT, Inches(6.55),
                                 T.CONTENT_WIDTH, Inches(0.5),
                                 note, size=T.PT_CAPTION, color=T.GRAY_80,
                                 line_spacing=1.3)
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Custom: Sequence diagram

def make_sequence_diagram(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    actors = data.get("actors", [])
    n = len(actors)
    seq_top = Inches(1.9)
    seq_left = T.CONTENT_LEFT
    seq_w = T.CONTENT_WIDTH
    actor_h = Inches(0.55)
    lifeline_h = Inches(3.7)
    gap = seq_w / n
    for i, name in enumerate(actors):
        cx_left = int(seq_left + gap * i + (gap - Inches(2.0)) / 2)
        S.add_filled_text_rect(slide, cx_left, seq_top, Inches(2.0), actor_h,
                               name, fill=T.BLUE_10, line=T.DEEP_BLUE,
                               text_color=T.DEEP_BLUE,
                               size=T.PT_BODY_SM, bold=True, rounded=True,
                               line_width=Pt(0.75))
        line_x = int(seq_left + gap * i + gap / 2)
        S.add_thin_line(slide, line_x, int(seq_top + actor_h),
                        0, color=T.GRAY_20, weight=Pt(1))
        slide.shapes.add_connector(
            1, line_x, int(seq_top + actor_h), line_x,
            int(seq_top + actor_h + lifeline_h)
        ).line.color.rgb = T.GRAY_20
    msgs = data.get("messages", [])
    msg_top = int(seq_top + actor_h + Inches(0.4))
    msg_gap = Inches(0.55)
    for i, msg in enumerate(msgs):
        S.add_textbox(slide, seq_left, msg_top + msg_gap * i,
                      seq_w, Inches(0.4),
                      f"{i+1}.  {msg}",
                      size=T.PT_BODY_SM, color=T.BLACK)
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Custom: ERD

def make_erd(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    entities = data.get("entities", [])
    n = len(entities)
    if n == 0:
        _attach_note(slide, data)
        _add_footer(slide, data["no"])
        return slide
    box_w = Inches(3.4)
    box_h = Inches(3.6)
    gap = (T.CONTENT_WIDTH - box_w * n) / (n + 1)
    top = Inches(2.0)
    centers = []
    for i, ent in enumerate(entities):
        left = T.CONTENT_LEFT + gap + (box_w + gap) * i
        S.add_filled_text_rect(slide, int(left), top, box_w, Inches(0.55),
                               ent["name"], fill=T.DEEP_BLUE, line=None,
                               text_color=T.WHITE, size=T.PT_BODY_SM,
                               bold=True, align=PP_ALIGN.LEFT)
        S.add_rect(slide, int(left), top + Inches(0.55), box_w,
                   box_h - Inches(0.55), fill=T.WHITE, line=T.DEEP_BLUE,
                   line_width=Pt(1.0))
        pad = Inches(0.2)
        S.add_multi_line_textbox(
            slide, int(left) + pad, top + Inches(0.55) + pad,
            box_w - 2 * pad, box_h - Inches(0.55) - 2 * pad,
            ent.get("fields", []),
            size=T.PT_BODY_SM, color=T.BLACK, line_spacing=1.45,
            name=T.FONT_CODE,
        )
        centers.append((int(left + box_w / 2), int(top + box_h / 2)))
    relations = data.get("relations", [])
    for (i, j, label) in relations:
        x1, _y1 = centers[i]
        x2, _y2 = centers[j]
        x1_edge = int(x1 + box_w / 2 if i < j else x1 - box_w / 2)
        x2_edge = int(x2 - box_w / 2 if i < j else x2 + box_w / 2)
        y_mid = int(top + box_h / 2)
        conn = slide.shapes.add_connector(1, x1_edge, y_mid, x2_edge, y_mid)
        conn.line.color.rgb = T.DEEP_BLUE
        conn.line.width = Pt(1.5)
        S.add_textbox(slide, min(x1_edge, x2_edge),
                      y_mid - Inches(0.35),
                      abs(x2_edge - x1_edge), Inches(0.3),
                      label, size=T.PT_CAPTION, color=T.DEEP_BLUE,
                      align=PP_ALIGN.CENTER, bold=True)
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Custom: 4 cards (personas / role split)

def make_card_grid(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    cards = data.get("cards", [])
    cols = data.get("card_cols", 2)
    rows = (len(cards) + cols - 1) // cols
    gap = Inches(0.25)
    body = data.get("body", [])
    bottom_reserve = Inches(0.7) if body else Inches(0)
    avail_h = Inches(4.6) - gap * (rows - 1) - bottom_reserve
    avail_w = T.CONTENT_WIDTH - gap * (cols - 1)
    card_w = avail_w / cols
    card_h = avail_h / rows
    top0 = Inches(1.85)
    for i, c in enumerate(cards):
        r = i // cols
        col = i % cols
        left = T.CONTENT_LEFT + (card_w + gap) * col
        top = top0 + (card_h + gap) * r
        S.add_card(slide, int(left), int(top), int(card_w), int(card_h),
                   c["title"], c.get("body", []),
                   accent=c.get("accent", False))
    if body:
        S.add_multi_line_textbox(slide, T.CONTENT_LEFT, Inches(6.65),
                                 T.CONTENT_WIDTH, Inches(0.5),
                                 body, size=T.PT_CAPTION, color=T.GRAY_80,
                                 line_spacing=1.3)
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Custom: badge row (formula)

def make_badge_row(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    badges = data.get("badges", [])
    n = len(badges)
    if n:
        gap = Inches(0.18)
        total_w = T.CONTENT_WIDTH
        plus_w = Inches(0.4)
        n_plus = n - 1
        bw = (total_w - gap * (n - 1) - plus_w * n_plus) / n
        cur = T.CONTENT_LEFT
        top = Inches(2.4)
        for i, label in enumerate(badges):
            S.add_filled_text_rect(slide, int(cur), top, int(bw),
                                   Inches(1.0), label,
                                   fill=T.BLUE_10, line=T.DEEP_BLUE,
                                   text_color=T.DEEP_BLUE, size=T.PT_BODY,
                                   bold=True, rounded=True,
                                   line_width=Pt(1.0))
            cur += bw + gap
            if i < n - 1:
                S.add_textbox(slide, int(cur), top + Inches(0.25),
                              int(plus_w), Inches(0.5), "+",
                              size=Pt(28), bold=True, color=T.DEEP_BLUE,
                              align=PP_ALIGN.CENTER,
                              anchor=MSO_ANCHOR.MIDDLE)
                cur += plus_w + gap
    body = data.get("body", [])
    if body:
        S.add_multi_line_textbox(slide, T.CONTENT_LEFT, Inches(4.0),
                                 T.CONTENT_WIDTH, Inches(2.6),
                                 body, size=T.PT_BODY, color=T.BLACK,
                                 line_spacing=1.55, bullet=True)
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Custom: comparison

def make_comparison(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    columns = data.get("columns", [])
    n = len(columns)
    if n:
        gap = Inches(0.3)
        col_w = (T.CONTENT_WIDTH - gap * (n - 1)) / n
        top = Inches(2.0)
        for i, col in enumerate(columns):
            left = T.CONTENT_LEFT + (col_w + gap) * i
            accent = col.get("accent", False)
            S.add_filled_text_rect(
                slide, int(left), top, int(col_w), Inches(0.6),
                col["title"],
                fill=T.DEEP_BLUE if accent else T.WHITE,
                line=None if accent else T.GRAY_20,
                text_color=T.WHITE if accent else T.BLACK,
                size=T.PT_BODY, bold=True, align=PP_ALIGN.CENTER)
            S.add_rect(slide, int(left), int(top + Inches(0.6)),
                       int(col_w), Inches(4.4),
                       fill=T.BLUE_10 if accent else T.GRAY_5,
                       line=T.DEEP_BLUE if accent else T.GRAY_20,
                       line_width=Pt(0.75))
            pad = Inches(0.2)
            S.add_multi_line_textbox(
                slide, int(left) + pad, int(top + Inches(0.6)) + pad,
                int(col_w) - 2 * pad, Inches(4.0),
                col.get("body", []),
                size=T.PT_BODY_SM, color=T.BLACK,
                line_spacing=1.45, bullet=True)
    body = data.get("body", [])
    if body:
        S.add_multi_line_textbox(slide, T.CONTENT_LEFT, Inches(6.7),
                                 T.CONTENT_WIDTH, Inches(0.5),
                                 body, size=T.PT_CAPTION, color=T.GRAY_80,
                                 line_spacing=1.3)
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Custom: API table

def make_api_table(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    rows = data.get("rows", [])
    headers = data.get("headers", ["Method", "Path", "Description"])
    S.add_table(slide, T.CONTENT_LEFT, Inches(1.95),
                T.CONTENT_WIDTH, Inches(4.5),
                headers, rows,
                col_widths=data.get("col_widths"))
    body = data.get("body", [])
    if body:
        S.add_multi_line_textbox(slide, T.CONTENT_LEFT, Inches(6.55),
                                 T.CONTENT_WIDTH, Inches(0.5),
                                 body, size=T.PT_CAPTION, color=T.GRAY_80,
                                 line_spacing=1.3)
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Cycle

def make_cycle(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    items = data.get("cycle_items", [])
    n = len(items)
    if n:
        cx = int(T.SLIDE_W / 2)
        cy = Inches(4.4)
        radius_x = Inches(2.4)
        radius_y = Inches(1.8)
        import math
        box_w = Inches(2.2)
        box_h = Inches(0.9)
        for i, label in enumerate(items):
            theta = -math.pi / 2 + 2 * math.pi * i / n
            x = cx + math.cos(theta) * radius_x - box_w / 2
            y = cy + math.sin(theta) * radius_y - box_h / 2
            S.add_filled_text_rect(slide, int(x), int(y), box_w, box_h,
                                   label, fill=T.BLUE_10, line=T.DEEP_BLUE,
                                   text_color=T.DEEP_BLUE, size=T.PT_BODY,
                                   bold=True, rounded=True,
                                   line_width=Pt(1.0))
    body = data.get("body", [])
    if body:
        S.add_multi_line_textbox(slide, T.CONTENT_LEFT, Inches(6.5),
                                 T.CONTENT_WIDTH, Inches(0.7),
                                 body, size=T.PT_CAPTION, color=T.GRAY_80,
                                 line_spacing=1.4)
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide


# ---------------------------------------------------------------- Tree (folder)

def make_tree(prs, data):
    slide = _add_blank_slide(prs)
    _add_title(slide, data["title"])
    S.add_thin_line(slide, T.CONTENT_LEFT, Inches(1.4), Inches(0.8),
                    color=T.DEEP_BLUE, weight=Pt(2.5))
    body = data.get("body", [])
    S.add_multi_line_textbox(
        slide, T.CONTENT_LEFT, Inches(1.95),
        Inches(5.0), Inches(5.0),
        body, size=T.PT_BODY, color=T.BLACK,
        line_spacing=1.55, bullet=True,
    )
    tree_left = T.CONTENT_LEFT + Inches(5.6)
    tree_w = T.CONTENT_WIDTH - Inches(5.6)
    S.add_code_block(slide, tree_left, Inches(1.95), tree_w, Inches(5.0),
                     data.get("tree_lines", []))
    _attach_note(slide, data)
    _add_footer(slide, data["no"])
    return slide
