"""Shape helpers — boxes, arrows, badges, lists, tables."""

from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

import theme as T


def set_run_font(run, size=None, bold=False, color=None, name=None, italic=False):
    f = run.font
    f.name = name or T.FONT_PRIMARY
    if size is not None:
        f.size = size
    f.bold = bold
    f.italic = italic
    if color is not None:
        f.color.rgb = color


def add_textbox(slide, left, top, width, height, text, *,
                size=None, bold=False, color=None, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, name=None, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size or T.PT_BODY, bold=bold,
                 color=color or T.BLACK, name=name)
    return tb


def add_multi_line_textbox(slide, left, top, width, height, lines, *,
                           size=None, bold=False, color=None,
                           align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                           bullet=False, line_spacing=1.3, name=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = (("• " if bullet else "") + line)
        set_run_font(run, size=size or T.PT_BODY, bold=bold,
                     color=color or T.BLACK, name=name)
    return tb


def add_rect(slide, left, top, width, height, *,
             fill=None, line=None, line_width=Pt(0.75),
             rounded=False, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = line_width
    if not shadow:
        try:
            s.shadow.inherit = False
        except Exception:
            pass
    if rounded:
        try:
            s.adjustments[0] = 0.08
        except Exception:
            pass
    return s


def add_filled_text_rect(slide, left, top, width, height, text, *,
                         fill=None, line=None, text_color=None,
                         size=None, bold=False, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE, rounded=False,
                         line_width=Pt(0.75), name=None):
    rect = add_rect(slide, left, top, width, height,
                    fill=fill, line=line, line_width=line_width, rounded=rounded)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(60000)
    tf.margin_right = Emu(60000)
    tf.margin_top = Emu(40000)
    tf.margin_bottom = Emu(40000)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size or T.PT_BODY, bold=bold,
                 color=text_color or T.BLACK, name=name)
    return rect


def add_callout_box(slide, left, top, width, height, text, *,
                    size=None, bold=False, align=PP_ALIGN.LEFT):
    rect = add_filled_text_rect(slide, left, top, width, height, text,
                                fill=T.BLUE_10, line=None,
                                text_color=T.BLACK, size=size or T.PT_BODY,
                                bold=bold, align=align,
                                anchor=MSO_ANCHOR.MIDDLE)
    bar = add_rect(slide, left, top, Emu(45720), height, fill=T.DEEP_BLUE)
    return rect, bar


def add_badge(slide, left, top, width, height, text, *,
              size=None, name=None):
    return add_filled_text_rect(slide, left, top, width, height, text,
                                fill=T.BLUE_10, line=None,
                                text_color=T.DEEP_BLUE, size=size or T.PT_BADGE,
                                bold=True, align=PP_ALIGN.CENTER, rounded=True,
                                anchor=MSO_ANCHOR.MIDDLE, name=name)


def add_arrow(slide, left, top, width, height, color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color or T.DEEP_BLUE
    s.line.fill.background()
    return s


def add_horizontal_flow(slide, left, top, total_width, height, items, *,
                        emphasize_index=None, gap_ratio=0.18):
    n = len(items)
    gap = total_width * gap_ratio / max(n - 1, 1)
    box_w = (total_width - gap * (n - 1)) / n
    arrow_w = gap * 0.85
    arrow_h = height * 0.35
    cur_left = left
    for i, label in enumerate(items):
        is_emph = (emphasize_index is None) or (i == emphasize_index)
        fill = T.BLUE_10 if is_emph else None
        line = T.DEEP_BLUE if is_emph else T.GRAY_20
        text_color = T.DEEP_BLUE if is_emph else T.BLACK
        bold = is_emph
        add_filled_text_rect(slide, int(cur_left), top, int(box_w), height, label,
                             fill=fill, line=line, text_color=text_color,
                             size=T.PT_BODY, bold=bold, rounded=True,
                             line_width=Pt(1.0))
        if i < n - 1:
            ax = int(cur_left + box_w + (gap - arrow_w) / 2)
            ay = int(top + (height - arrow_h) / 2)
            add_arrow(slide, ax, ay, int(arrow_w), int(arrow_h),
                      color=T.DEEP_BLUE)
        cur_left += box_w + gap


def add_card(slide, left, top, width, height, title, body_lines, *,
             accent=False):
    fill = T.BLUE_10 if accent else T.WHITE
    line = T.DEEP_BLUE if accent else T.GRAY_20
    rect = add_rect(slide, left, top, width, height,
                    fill=fill, line=line, line_width=Pt(1.0), rounded=True)
    pad = Emu(120000)
    add_textbox(slide, left + pad, top + pad,
                width - 2 * pad, Emu(450000), title,
                size=T.PT_BODY, bold=True,
                color=T.DEEP_BLUE if accent else T.BLACK)
    if body_lines:
        add_multi_line_textbox(
            slide, left + pad, top + pad + Emu(420000),
            width - 2 * pad, height - 2 * pad - Emu(420000),
            body_lines, size=T.PT_BODY_SM, color=T.GRAY_80,
            line_spacing=1.25)
    return rect


def add_placeholder_box(slide, left, top, width, height, caption):
    rect = add_rect(slide, left, top, width, height,
                    fill=T.GRAY_5, line=T.GRAY_20, line_width=Pt(0.75),
                    rounded=False)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = caption
    set_run_font(run, size=T.PT_BODY_SM, color=T.GRAY_80,
                 name=T.FONT_PRIMARY, italic=True)
    return rect


def add_code_block(slide, left, top, width, height, lines, *,
                   size=None, highlight_indices=None):
    rect = add_rect(slide, left, top, width, height,
                    fill=T.GRAY_5, line=T.GRAY_20, line_width=Pt(0.5))
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(150000)
    tf.margin_right = Emu(150000)
    tf.margin_top = Emu(100000)
    tf.margin_bottom = Emu(100000)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    highlights = highlight_indices or set()
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.3
        run = p.add_run()
        run.text = line if line else " "
        is_highlight = i in highlights
        set_run_font(run, size=size or T.PT_CODE,
                     color=T.DEEP_BLUE if is_highlight else T.BLACK,
                     bold=is_highlight, name=T.FONT_CODE)
    return rect


def add_thin_line(slide, left, top, width, color=None, weight=Pt(1.5)):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color or T.DEEP_BLUE
    line.line.width = weight
    return line


def add_table(slide, left, top, width, height, headers, rows, *,
              col_widths=None):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            if w is not None:
                table.columns[i].width = w
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = T.WHITE
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Emu(80000)
        tf.margin_right = Emu(80000)
        tf.margin_top = Emu(40000)
        tf.margin_bottom = Emu(40000)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = h
        set_run_font(run, size=T.PT_BODY_SM, bold=True, color=T.BLACK)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = T.WHITE
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Emu(80000)
            tf.margin_right = Emu(80000)
            tf.margin_top = Emu(40000)
            tf.margin_bottom = Emu(40000)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            set_run_font(run, size=T.PT_BODY_SM, color=T.BLACK)
    return table


def add_check_list(slide, left, top, width, height, items, *,
                   size=None, line_spacing=1.5):
    lines = [f"☐  {it}" for it in items]
    return add_multi_line_textbox(slide, left, top, width, height, lines,
                                  size=size or T.PT_BODY,
                                  line_spacing=line_spacing)


def set_slide_notes(slide, note_text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = note_text or ""
    for p in tf.paragraphs:
        for r in p.runs:
            set_run_font(r, size=Pt(11), color=T.BLACK, name=T.FONT_PRIMARY)
